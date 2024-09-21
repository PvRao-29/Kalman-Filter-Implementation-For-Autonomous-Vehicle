import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation(output_file='Kalman_Filter_Position_Estimation.pptx'):
    # Initialize Presentation
    prs = Presentation()

    # Define Slide Layouts
    title_slide_layout = prs.slide_layouts[0]  # Title Slide
    content_slide_layout = prs.slide_layouts[1]  # Title and Content
    section_header_layout = prs.slide_layouts[2]  # Section Header
    two_content_layout = prs.slide_layouts[3]  # Two Content
    comparison_layout = prs.slide_layouts[4]  # Comparison
    blank_layout = prs.slide_layouts[6]  # Blank

    # Slide 1: Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Kalman Filter-Based Position Estimation for Vehicle Movement"
    subtitle.text = "Integrating GPS, IMU, and Wheel Speed Data\nYour Name\nDate"

    # Optional: Add Logo
    logo_path = 'logo.png'  # Ensure this file exists or comment out the following lines
    if os.path.exists(logo_path):
        left = Inches(9)  # Positioning the logo on the right
        top = Inches(0.5)
        height = Inches(1)
        slide.shapes.add_picture(logo_path, left, top, height=height)

    # Slide 2: Project Overview
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Project Overview"

    content = slide.placeholders[1]
    text = (
        "• **Objective:**\n"
        "  - Accurately estimate a vehicle's position by fusing data from GPS, IMU, and wheel speed sensors using a Kalman Filter.\n"
        "• **Key Components:**\n"
        "  - Data Collection\n"
        "  - Noise Calibration\n"
        "  - Kalman Filter Implementation\n"
        "  - Position Estimation and Visualization"
    )
    content.text = text

    # Slide 3: Data Sources
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Data Sources"

    content = slide.placeholders[1]
    text = (
        "• **Stationary Data (`stationary.csv`):**\n"
        "  - Purpose: Calibration of GPS noise.\n"
        "  - Key Fields: Latitude, Longitude.\n\n"
        "• **Movement Data (`movement.csv`):**\n"
        "  - Purpose: Real-time position estimation.\n"
        "  - Key Fields: Time (Seconds, Nanoseconds), Latitude, Longitude, LinearAccel.x/y/z, LeftFrontSpeed, RightFrontSpeed, LeftBackSpeed, RightBackSpeed."
    )
    content.text = text

    # Slide 4: Methodology - Kalman Filter
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Methodology - Kalman Filter"

    content = slide.placeholders[1]
    text = (
        "• **What is a Kalman Filter?**\n"
        "  - A recursive algorithm for estimating the state of a dynamic system from noisy measurements.\n\n"
        "• **Why Use It?**\n"
        "  - Combines predictions from a mathematical model with real-time measurements to produce optimal estimates.\n\n"
        "• **Components:**\n"
        "  - **State Vector:** `[Latitude, Longitude, Velocity_Lat, Velocity_Lon]`\n"
        "  - **Covariance Matrices:** `P` (State Covariance), `R` (Measurement Noise), `Q` (Process Noise)\n"
        "  - **Measurement Matrix:** `H`"
    )
    content.text = text

    # Slide 5: Implementation Details
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Implementation Details"

    content = slide.placeholders[1]
    text = (
        "• **Initialization:**\n"
        "  - State vector initialized with the first valid GPS measurement.\n"
        "  - Initial covariance matrices set based on calibrated noise.\n\n"
        "• **Prediction Step:**\n"
        "  - State transition matrix `F` updated with time step `delta_t`.\n"
        "  - State and covariance predictions.\n\n"
        "• **Update Step:**\n"
        "  - Incorporation of GPS measurements when available.\n"
        "  - Use of IMU and wheel speed data when GPS is unavailable.\n\n"
        "• **Handling Missing Data:**\n"
        "  - Strategies for dealing with NaNs and incomplete sensor readings."
    )
    content.text = text

    # Slide 6: Code Structure
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Code Structure"

    content = slide.placeholders[1]
    text = (
        "• **Modules and Functions:**\n"
        "  - Data Loading and Preprocessing\n"
        "  - Noise Calibration\n"
        "  - Kalman Filter Operations\n"
        "  - Position Estimation and Storage\n\n"
        "• **Key Libraries:**\n"
        "  - `numpy`, `pandas` for data manipulation.\n"
        "  - `python-pptx` for presentation generation."
    )
    content.text = text

    # Slide 7: Results - Position Estimates
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Results - Position Estimates"

    content = slide.placeholders[1]
    text = (
        "• **Overview:**\n"
        "  - Generated `position.csv` containing refined latitude and longitude estimates.\n\n"
        "• **Sample Data:**\n"
        "  | Latitude | Longitude |\n"
        "  |----------|-----------|\n"
        "  | 34.05    | -118.25   |\n"
        "  | 34.051   | -118.251  |\n"
        "  | ...      | ...       |\n"
        "  | 34.075   | -118.275  |"
    )
    content.text = text

    # Slide 8: Visualization
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Visualization"

    content = slide.placeholders[1]
    text = (
        "• **Scatter Plot of Position Estimates:**\n"
        "  - All position points plotted.\n"
        "  - Start and end points highlighted.\n"
        "  - Intermediate points optionally marked.\n\n"
        "• **Color Mapping:**\n"
        "  - Progression over time indicated via color gradients."
    )
    content.text = text

    # Slide 9: Performance Enhancements
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Performance Enhancements"

    content = slide.placeholders[1]
    text = (
        "• **Optimizations Implemented:**\n"
        "  - State Initialization with First GPS Measurement.\n"
        "  - Accurate Time Step (`delta_t`) Calculation.\n"
        "  - Efficient Data Handling and Batch Writing.\n"
        "  - Robust Error Handling and Logging.\n\n"
        "• **Potential Further Improvements:**\n"
        "  - Implementing Extended Kalman Filter (EKF) for Nonlinear Dynamics.\n"
        "  - Integrating Additional Sensors for Enhanced Accuracy."
    )
    content.text = text

    # Slide 10: Conclusion
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Conclusion"

    content = slide.placeholders[1]
    text = (
        "• **Achievements:**\n"
        "  - Successfully implemented a Kalman Filter for position estimation.\n"
        "  - Integrated multiple data sources to enhance accuracy.\n"
        "  - Developed visualization tools to interpret results.\n\n"
        "• **Impact:**\n"
        "  - Improved reliability of vehicle tracking systems.\n"
        "  - Foundation for real-time navigation and autonomous driving applications."
    )
    content.text = text

    # Slide 11: Future Work
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Future Work"

    content = slide.placeholders[1]
    text = (
        "• **Enhancements:**\n"
        "  - Incorporate Heading Information for Directional Accuracy.\n"
        "  - Utilize FilterPy or Other Libraries for Advanced Filtering Techniques.\n"
        "  - Develop a User Interface for Real-Time Data Querying and Visualization.\n\n"
        "• **Applications:**\n"
        "  - Autonomous Vehicles\n"
        "  - Fleet Management Systems\n"
        "  - Enhanced Navigation Tools"
    )
    content.text = text

    # Slide 12: Q&A
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Questions & Answers"

    content = slide.placeholders[1]
    text = "Feel free to ask any questions or seek clarifications regarding the project."
    content.text = text

    # Save Presentation
    prs.save(output_file)
    print(f"Presentation '{output_file}' created successfully.")

if __name__ == "__main__":
    create_presentation()
